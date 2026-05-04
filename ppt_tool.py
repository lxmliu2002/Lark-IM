import os
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
CAPTURE_WIDTH = 1600
CAPTURE_HEIGHT = 900


def render_html_slides_to_ppt(slide_files: list[str], output_path: str) -> str:
    if not slide_files:
        raise ValueError("No HTML slide files were provided for PPT export.")

    browser_path = _find_browser()
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)

    image_dir = output_file.parent / f"{output_file.stem}_images"
    image_dir.mkdir(parents=True, exist_ok=True)
    image_paths: list[Path] = []
    for index, slide_file in enumerate(slide_files, start=1):
        image_path = image_dir / f"slide_{index:02}.png"
        _capture_slide(browser_path, Path(slide_file), image_path)
        image_paths.append(image_path)

    _images_to_ppt(image_paths, output_file)
    return str(output_file)


def _capture_slide(browser_path: str, slide_file: Path, image_path: Path) -> None:
    slide_uri = slide_file.resolve().as_uri()
    command = [
        browser_path,
        "--headless",
        "--disable-gpu",
        "--hide-scrollbars",
        "--window-size=%d,%d" % (CAPTURE_WIDTH, CAPTURE_HEIGHT),
        "--screenshot=%s" % str(image_path.resolve()),
        slide_uri,
    ]

    result = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=45,
        check=False,
    )

    if result.returncode != 0 or not image_path.exists():
        message = result.stderr.strip() or result.stdout.strip() or "Unknown browser capture error."
        raise RuntimeError(f"Failed to render HTML slide to image: {message}")


def _images_to_ppt(image_paths: list[Path], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    blank_layout = prs.slide_layouts[6]
    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            Inches(0),
            Inches(0),
            width=Inches(SLIDE_WIDTH),
            height=Inches(SLIDE_HEIGHT),
        )

    prs.save(str(output_path))


def _find_browser() -> str:
    env_path = os.getenv("HTML_RENDER_BROWSER")
    if env_path and Path(env_path).exists():
        return env_path

    candidates = [
        Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
        Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
    ]

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    raise RuntimeError(
        "No supported browser was found. Please install Microsoft Edge or Chrome, "
        "or set HTML_RENDER_BROWSER in .env to the browser executable path."
    )
